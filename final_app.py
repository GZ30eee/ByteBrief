import os
import streamlit as st
st.set_page_config(page_title="Smart Text Summarizer", layout="wide")
import requests
import spacy
import docx
import PyPDF2
from pptx import Presentation
from io import BytesIO
import fitz  # PyMuPDF for PDF rendering
import numpy as np
from youtube_transcript_api import YouTubeTranscriptApi
import whisper
from transformers import pipeline
from keybert import KeyBERT
from gensim import corpora
from gensim.models import LdaModel
import matplotlib.pyplot as plt
from wordcloud import WordCloud
import plotly.express as px
import pandas as pd
import re
import nltk
from nltk.tokenize import sent_tokenize
import easyocr
from PIL import Image
import tempfile
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

# Initialize models
@st.cache_resource
def load_models():
    nlp = spacy.load("en_core_web_sm")
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
    keyword_model = KeyBERT(model='paraphrase-MiniLM-L6-v2')
    sentiment_analyzer = pipeline("sentiment-analysis")
    whisper_model = whisper.load_model("tiny")
    reader = None
    try:
        reader = easyocr.Reader(['hi', 'en'])  # Support Hindi and English
    except Exception as e:
        st.warning(f"Failed to initialize EasyOCR with Hindi: {str(e)}. Falling back to English-only OCR.")
        try:
            reader = easyocr.Reader(['en'])
        except Exception as e2:
            st.error(f"Failed to initialize EasyOCR for English: {str(e2)}. Image OCR will not be available.")
    return nlp, summarizer, keyword_model, sentiment_analyzer, whisper_model, reader

nlp, summarizer, keyword_model, sentiment_analyzer, whisper_model, ocr_reader = load_models()

# Word count function
def count_words(text):
    return len([word for word in text.split() if word.strip()])

# Function to fetch random text with at least 30 words
def fetch_random_paragraph():
    try:
        max_attempts = 5
        for _ in range(max_attempts):
            response = requests.get("https://en.wikipedia.org/api/rest_v1/page/random/summary")
            if response.status_code == 200:
                text = response.json().get("extract", "No content available.")
                if count_words(text) >= 30:
                    return text
        return "Could not fetch a paragraph with at least 30 words. Please try again."
    except requests.exceptions.RequestException:
        return "Connection error. Please check your internet connection."

# Transformer-based summarizer
def transformer_summarizer(text, target_word_count, keywords=None):
    try:
        if keywords:
            keyword_text = " ".join(keywords) + "."
            text = text + " Keywords: " + keyword_text
        tokens = summarizer.tokenizer(text, truncation=False, return_tensors="pt")["input_ids"][0]
        max_input_tokens = 1024 - 50
        if len(tokens) > max_input_tokens:
            truncated_tokens = tokens[:max_input_tokens]
            text = summarizer.tokenizer.decode(truncated_tokens, skip_special_tokens=True)
        
        max_length = min(target_word_count + 20, 1024)
        min_length = max(target_word_count // 2, 10)
        summary = summarizer(text, max_length=max_length, min_length=min_length, do_sample=False)[0]['summary_text']
        
        sentences = sent_tokenize(summary)
        final_summary = []
        current_length = 0
        for sent in sentences:
            sent_words = count_words(sent)
            if current_length + sent_words <= target_word_count:
                final_summary.append(sent)
                current_length += sent_words
            else:
                break
        final_text = ' '.join(final_summary)
        return final_text, len(final_summary), count_words(final_text)
    except Exception as e:
        return f"Error in summarization: {str(e)}", 0, 0

# Convert summary to bullet points
def convert_to_bullets(summary_text):
    sentences = sent_tokenize(summary_text)
    return '\n\n'.join([f"• {sent}" for sent in sentences])

# Extract text from uploaded files
def extract_text_from_uploaded_file(uploaded_file):
    try:
        file_extension = uploaded_file.name.split('.')[-1].lower()
        text = ""
        if file_extension == 'txt':
            text = uploaded_file.read().decode("utf-8")
        elif file_extension == 'pdf':
            reader = PyPDF2.PdfReader(uploaded_file)
            text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif file_extension == 'docx':
            doc = docx.Document(uploaded_file)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif file_extension in ['ppt', 'pptx']:
            prs = Presentation(uploaded_file)
            text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])
        elif file_extension in ['jpg', 'jpeg', 'png']:
            if ocr_reader is None:
                st.error("EasyOCR is not available. Cannot process image files.")
                return ""
            image = Image.open(uploaded_file)
            image_array = np.array(image)  # Convert PIL Image to NumPy array
            result = ocr_reader.readtext(image_array, detail=0)
            text = " ".join(result)
        elif file_extension in ['mp3', 'wav']:
            with st.spinner("Transcribing audio..."):
                result = whisper_model.transcribe(uploaded_file.name)
                text = result["text"]
        if count_words(text) < 30:
            st.markdown(f"<p style='color:red;'>Error: Uploaded file contains {count_words(text)} words, fewer than the required 30 words.</p>", unsafe_allow_html=True)
            return ""
        return text
    except Exception as e:
        st.error(f"Error processing file: {str(e)}")
        return ""

# YouTube transcript extraction
def extract_youtube_transcript(youtube_url, max_retries=3):
    try:
        from pytube import YouTube
        from youtube_transcript_api import YouTubeTranscriptApi
        from youtube_transcript_api._errors import TranscriptsDisabled, NoTranscriptFound
        
        video_id_match = re.findall(r"(?:v=|\/)([0-9A-Za-z_-]{11}).*", youtube_url)
        if not video_id_match:
            return "Invalid YouTube URL. Please provide a valid video link."
        
        video_id = video_id_match[0]
        st.info(f"Processing YouTube video ID: {video_id}")

        # Try YouTubeTranscriptApi first
        try:
            transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
            transcript = None
            
            # Try English first, then Hindi, then first available
            for lang in ['en', 'hi']:
                try:
                    transcript = transcript_list.find_transcript([lang])
                    break
                except:
                    continue
            
            if not transcript:
                transcript = transcript_list.find_transcript([t.language_code for t in transcript_list])
            
            transcript_data = transcript.fetch()
            text = " ".join([entry['text'] for entry in transcript_data])
            text = text.strip()
            
            if count_words(text) >= 30:
                return text
            else:
                st.warning(f"Transcript contains {count_words(text)} words, fewer than the required 30 words. Trying fallback methods...")
                
        except (TranscriptsDisabled, NoTranscriptFound) as e:
            st.warning("Transcripts disabled or not available via API. Trying fallback methods...")
        except Exception as e:
            st.warning(f"Error with YouTubeTranscriptApi: {str(e)}. Trying fallback methods...")

        # Fallback to pytube captions
        try:
            yt = YouTube(youtube_url)
            captions = yt.captions
            
            if captions:
                caption = None
                for code in ['en', 'hi']:
                    try:
                        caption = captions[code] if code in captions else None
                        if caption: break
                    except:
                        continue
                
                if not caption and captions:
                    caption = next(iter(captions.values()))
                
                if caption:
                    text = caption.generate_srt_captions()
                    text = re.sub(r'\d+\n\d{2}:\d{2}:\d{2},\d{3} --> \d{2}:\d{2}:\d{2},\d{3}\n', '', text)
                    text = re.sub(r'\n+', ' ', text).strip()
                    
                    if count_words(text) >= 30:
                        return text
                    else:
                        st.warning(f"Captions contain {count_words(text)} words, fewer than the required 30 words.")
            
        except Exception as e:
            st.warning(f"Error with pytube captions: {str(e)}")

        # Final fallback to Whisper
        try:
            st.info("Transcribing audio with Whisper...")
            yt = YouTube(youtube_url)
            audio_stream = yt.streams.filter(only_audio=True).order_by('abr').last()
            
            if not audio_stream:
                return "No audio stream available for this video."
            
            with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as temp_audio:
                audio_file = temp_audio.name
                try:
                    audio_stream.download(filename=audio_file)
                    
                    if os.path.getsize(audio_file) == 0:
                        return "Downloaded audio file is empty."
                    
                    result = whisper_model.transcribe(audio_file)
                    text = result["text"].strip()
                    
                    if count_words(text) >= 30:
                        return text
                    else:
                        st.warning(f"Whisper transcription contains {count_words(text)} words, fewer than the required 30 words.")
                        return text if text else "Whisper transcription is empty."
                
                finally:
                    if os.path.exists(audio_file):
                        try:
                            os.remove(audio_file)
                        except:
                            pass
        
        except Exception as e:
            return f"Error processing YouTube video: {str(e)}"
        
        return "Could not extract transcript from this video."

    except Exception as e:
        return f"Error processing YouTube video: {str(e)}"
    
# Export summary as DOCX, PDF, or TXT
def export_summary(summary_text, format_type):
    if format_type == "DOCX":
        doc = docx.Document()
        doc.add_paragraph(summary_text)
        doc_io = BytesIO()
        doc.save(doc_io)
        doc_io.seek(0)
        return doc_io, "summary.docx"
    elif format_type == "PDF":
        pdf = fitz.open()
        page = pdf.new_page()
        page.insert_text((50, 50), summary_text)
        pdf_io = BytesIO()
        pdf.save(pdf_io)
        pdf_io.seek(0)
        return pdf_io, "summary.pdf"
    elif format_type == "TXT":
        txt_io = BytesIO()
        txt_io.write(summary_text.encode('utf-8'))
        txt_io.seek(0)
        return txt_io, "summary.txt"

# Topic modeling using LDA
def get_topics(text):
    try:
        doc = nlp(text)
        tokens = [token.lemma_ for token in doc if not token.is_stop and token.is_alpha]
        dictionary = corpora.Dictionary([tokens])
        corpus = [dictionary.doc2bow(tokens)]
        lda_model = LdaModel(corpus, num_topics=3, id2word=dictionary, passes=10)
        topics = lda_model.print_topics()
        return topics
    except:
        return []

# Generate word cloud
def generate_word_cloud(text):
    wordcloud = WordCloud(width=800, height=400, background_color='white').generate(text)
    plt.figure(figsize=(10, 5))
    plt.imshow(wordcloud, interpolation='bilinear')
    plt.axis('off')
    st.pyplot(plt)

# Streamlit UI
st.title("Smart Text Summarizer with AI/ML")

# Tabs for different input types
tab1, tab2, tab3 = st.tabs(["Text Input", "File Upload", "AI/ML Insights"])

with tab1:
    st.subheader("Enter Text")
    st.markdown("<p style='color:red;'>Paragraph must contain at least 30 words.</p>", unsafe_allow_html=True)
    if "text" not in st.session_state:
        st.session_state.text = ""
    text_area = st.text_area("Paste your text here", height=200, value=st.session_state.text)
    if count_words(text_area) < 30 and text_area.strip():
        st.markdown(f"<p style='color:red;'>Error: Paragraph contains {count_words(text_area)} words, fewer than the required 30 words.</p>", unsafe_allow_html=True)
    else:
        st.session_state.text = text_area
    if st.button("Get Random Text"):
        with st.spinner("Fetching random text..."):
            random_text = fetch_random_paragraph()
            random_word_count = count_words(random_text)
            if random_word_count >= 30:
                st.session_state.text = random_text
                st.rerun()
            else:
                st.markdown(f"<p style='color:red;'>Fetched paragraph contains {random_word_count} words, fewer than the required 30 words. Please try again.</p>", unsafe_allow_html=True)

with tab2:
    st.subheader("Upload File")
    st.markdown("<p style='color:red;'>Uploaded file must contain at least 30 words.</p>", unsafe_allow_html=True)
    uploaded_file = st.file_uploader("Upload text, image, or audio file", type=["txt", "pdf", "docx", "ppt", "pptx"])
    if uploaded_file:
        with st.spinner("Processing file..."):
            extracted_text = extract_text_from_uploaded_file(uploaded_file)
            if extracted_text and count_words(extracted_text) >= 30:
                st.session_state.text = extracted_text
                st.text_area("Extracted Text", value=st.session_state.text, height=200)
            elif extracted_text:
                st.markdown(f"<p style='color:red;'>Error: Uploaded file contains {count_words(extracted_text)} words, fewer than the required 30 words.</p>", unsafe_allow_html=True)

# with tab3:
#     st.subheader("YouTube Video")
#     youtube_url = st.text_input("Enter YouTube video URL")
#     if st.button("Extract Transcript"):
#         if youtube_url:
#             with st.spinner("Fetching transcript..."):
#                 transcript_text = extract_youtube_transcript(youtube_url)
#                 transcript_word_count = count_words(transcript_text)
#                 if transcript_text and transcript_word_count >= 30:
#                     st.session_state.text = transcript_text
#                     st.text_area("Extracted Transcript", value=st.session_state.text, height=200)
#                 elif transcript_text:
#                     st.session_state.text = transcript_text
#                     st.text_area("Extracted Transcript", value=st.session_state.text, height=200)
#                     st.markdown(f"<p style='color:red;'>Warning: Transcript contains {transcript_word_count} words, fewer than the required 30 words, but processing as audio transcription.</p>", unsafe_allow_html=True)

with tab3:
    st.subheader("AI/ML Insights")
    st.write("""
    This app leverages advanced AI/ML techniques:
    - **Summarization**: Uses BART (a transformer model) for abstractive summarization.
    - **Keyword Extraction**: Employs KeyBERT for context-aware keyword extraction.
    - **OCR**: EasyOCR for text extraction from images (English and Hindi only).
    - **Speech-to-Text**: Whisper for audio transcription.
    - **Topic Modeling**: LDA for identifying key topics.
    - **Sentiment Analysis**: Transformer-based sentiment classification.
    """)

# Main processing section
st.subheader("Summarization Options")
if st.session_state.text and count_words(st.session_state.text) >= 30:
    mode = st.selectbox("Select Output Format", ["Paragraph", "Bullet Points", "Custom Command"])
    input_word_count = count_words(st.session_state.text)
    target_words = st.slider("Target summary length (words)", min_value=30, max_value=input_word_count, value=min(150, input_word_count))
    export_format = st.selectbox("Export Format", ["DOCX", "PDF", "TXT"])

    # Keyword extraction
    if mode in ["Paragraph", "Bullet Points"]:
        keywords = keyword_model.extract_keywords(st.session_state.text, top_n=10)
        keyword_options = [kw[0] for kw in keywords]
        selected_keywords = st.multiselect("Select keywords to enhance summary", keyword_options, default=keyword_options[:3])
    else:
        selected_keywords = []

    # Custom command
    if mode == "Custom Command":
        custom_cmd = st.selectbox("Choose a Command", ["Give a Title", "Generate a Conclusion", "Make it Academic"])
    else:
        custom_cmd = None

    # Generate summary
    if st.button("Generate Summary"):
        with st.spinner("Generating summary..."):
            if mode == "Paragraph":
                result, num_sentences, num_words = transformer_summarizer(st.session_state.text, target_words, selected_keywords)
            elif mode == "Bullet Points":
                summary, num_sentences, num_words = transformer_summarizer(st.session_state.text, target_words, selected_keywords)
                result = convert_to_bullets(summary)
            elif mode == "Custom Command":
                doc = nlp(st.session_state.text)
                if custom_cmd == "Give a Title":
                    result = summarizer(st.session_state.text, max_length=15, min_length=5, do_sample=False)[0]['summary_text']
                    result = result.capitalize()
                    num_sentences = 1
                    num_words = count_words(result)
                elif custom_cmd == "Generate a Conclusion":
                    summary = summarizer(st.session_state.text, max_length=50, min_length=20, do_sample=False)[0]['summary_text']
                    result = f"In conclusion, {summary}"
                    num_sentences = len(sent_tokenize(result))
                    num_words = count_words(result)
                elif custom_cmd == "Make it Academic":
                    result = " ".join([token.text.upper() if token.pos_ in ["NOUN", "PROPN"] else token.text for token in doc])
                    sentences = sent_tokenize(result)
                    final_summary = []
                    current_length = 0
                    for sent in sentences:
                        sent_words = count_words(sent)
                        if current_length + sent_words <= target_words:
                            final_summary.append(sent)
                            current_length += sent_words
                        else:
                            break
                    result = " ".join(final_summary)
                    num_sentences = len(final_summary)
                    num_words = count_words(result)

            st.subheader(f"Summary - {num_sentences} Sentences, {num_words} Words")
            st.write(result)

            # Statistics
            with st.expander("Statistics"):
                word_count = count_words(st.session_state.text)
                characters = len(st.session_state.text)
                reduction = round((1 - (num_words / word_count)) * 100, 2) if word_count > 0 else 0
                max_sentiment_tokens = 512
                tokens = sentiment_analyzer.tokenizer(st.session_state.text, truncation=False, return_tensors="pt")["input_ids"][0]
                if len(tokens) > max_sentiment_tokens:
                    truncated_tokens = tokens[:max_sentiment_tokens]
                    truncated_text = sentiment_analyzer.tokenizer.decode(truncated_tokens, skip_special_tokens=True)
                else:
                    truncated_text = st.session_state.text
                sentiment = sentiment_analyzer(truncated_text)[0]
                st.write(f"Word count: {word_count}")
                st.write(f"Sentence count: {num_sentences}")
                st.write(f"Characters: {characters}")
                st.write(f"Reduction: {reduction}%")
                st.write(f"Sentiment: {sentiment['label']} (Score: {sentiment['score']:.2f})")

            # Visualizations
            with st.expander("Visualizations"):
                st.subheader("Word Cloud")
                generate_word_cloud(st.session_state.text)
                st.subheader("Topic Distribution")
                topics = get_topics(st.session_state.text)
                if topics:
                    topic_data = [{"Topic": f"Topic {i}", "Words": topic[1]} for i, topic in enumerate(topics)]
                    df = pd.DataFrame(topic_data)
                    fig = px.bar(df, x="Topic", y="Words", text="Words")
                    st.plotly_chart(fig)

            # Export
            st.subheader("Export Summary")
            doc_io, filename = export_summary(result, export_format)
            st.download_button(f"Download {export_format}", doc_io, filename)
else:
    st.markdown("<p style='color:red;'>Please provide a paragraph with at least 30 words to enable summarization.</p>", unsafe_allow_html=True)

# Clear button
if st.button("Clear All"):
    st.session_state.text = ""
    st.rerun()